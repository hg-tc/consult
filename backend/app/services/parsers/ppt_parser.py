"""
PPT 解析器（离线）
使用 python-pptx 提取文本/表格/备注；对幻灯片内图片调用离线 OCR；
输出 LlamaIndex Document 列表，按幻灯片元素分块。
"""

from __future__ import annotations

from typing import List, Dict, Any
from pathlib import Path
import tempfile
import os
import logging

logger = logging.getLogger(__name__)


def _try_import_python_pptx():
    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
        return Presentation, MSO_SHAPE_TYPE
    except Exception as e:
        raise ImportError("需要安装 python-pptx 才能解析 PPT 文件") from e


def _export_shape_image(shape, tmp_dir: Path) -> str | None:
    try:
        image = shape.image
        ext = image.ext or "png"
        out_path = tmp_dir / f"shape_{id(shape)}.{ext}"
        with open(out_path, "wb") as f:
            f.write(image.blob)
        return str(out_path)
    except Exception:
        return None


def _call_image_ocr(image_path: str) -> str:
    # 复用 image_parser 的 OCR
    from .image_parser import parse_image_to_documents  # type: ignore
    docs = parse_image_to_documents(image_path, source_metadata={"from_ppt": True})
    return docs[0].text if docs else ""


def _markdown_table_from_ppt_table(table) -> str:
    rows = table.rows
    cols = table.columns
    if len(cols) == 0 or len(rows) == 0:
        return ""
    # 读取第一行作为表头
    headers = []
    for c in range(len(cols)):
        cell = table.cell(0, c)
        headers.append((cell.text or "").strip())
    lines = ["| " + " | ".join(headers) + " |", "| " + " | ".join(["---"] * len(headers)) + " |"]

    for r in range(1, len(rows)):
        vals = []
        for c in range(len(cols)):
            cell = table.cell(r, c)
            vals.append((cell.text or "").strip())
        lines.append("| " + " | ".join(vals) + " |")
    return "\n".join(lines)


def parse_ppt_to_documents(file_path: str, source_metadata: Dict[str, Any] | None = None) -> List["Document"]:
    # 惰性导入 Document
    try:
        from llama_index.core import Document  # type: ignore
    except Exception:
        from llama_index import Document  # type: ignore

    Presentation, MSO_SHAPE_TYPE = _try_import_python_pptx()

    ppt_path = Path(file_path)
    if not ppt_path.exists():
        raise FileNotFoundError(f"PPT 文件不存在: {file_path}")

    prs = Presentation(str(ppt_path))
    logger.info(f"[PPT] 解析开始: file={ppt_path}, slides={len(prs.slides)}")

    base_meta: Dict[str, Any] = {
        "type": "ppt",
        "source_file": str(ppt_path),
        **(source_metadata or {}),
    }

    documents: List[Document] = []

    with tempfile.TemporaryDirectory() as td:
        tmp_dir = Path(td)
        for slide_idx, slide in enumerate(prs.slides, start=1):
            logger.info(f"[PPT] Slide {slide_idx}/{len(prs.slides)} 开始处理")
            slide_docs_before = len(documents)
            cnt_text = 0
            cnt_table = 0
            cnt_image = 0
            try:
                # 讲者备注
                try:
                    notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
                    if notes and notes.strip():
                        documents.append(Document(
                            text=f"# Slide {slide_idx} 备注\n\n{notes}",
                            metadata={**base_meta, "slide": slide_idx, "shape_type": "notes"}
                        ))
                except Exception as e:
                    logger.warning(f"[PPT] Slide {slide_idx} 备注解析失败: {e}")

                for shape in slide.shapes:
                    try:
                        stype = shape.shape_type
                    except Exception:
                        stype = None

                    # 文本框
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        text = shape.text or ""
                        if text.strip():
                            documents.append(Document(
                                text=f"# Slide {slide_idx} 文本\n\n{text}",
                                metadata={**base_meta, "slide": slide_idx, "shape_type": "text"}
                            ))
                            cnt_text += 1
                        continue

                    # 表格
                    if stype == MSO_SHAPE_TYPE.TABLE and hasattr(shape, "table"):
                        try:
                            md = _markdown_table_from_ppt_table(shape.table)
                            if md.strip():
                                documents.append(Document(
                                    text=f"# Slide {slide_idx} 表格\n\n{md}",
                                    metadata={**base_meta, "slide": slide_idx, "shape_type": "table"}
                                ))
                                cnt_table += 1
                        except Exception as e:
                            logger.warning(f"[PPT] Slide {slide_idx} 导出表格失败: {e}")
                        continue

                    # 图片（位图）
                    if hasattr(shape, "image"):
                        img_path = _export_shape_image(shape, tmp_dir)
                        if img_path:
                            logger.info(f"[PPT] Slide {slide_idx} 导出图片: {img_path}")
                            import time as _t
                            _t0 = _t.time()
                            try:
                                ocr_text = _call_image_ocr(img_path)
                                logger.info(f"[PPT] Slide {slide_idx} 图片 OCR 完成: took={_t.time()-_t0:.3f}s, has_text={bool(ocr_text.strip())}")
                            except Exception as e:
                                logger.warning(f"OCR 图片失败，退化为占位: {e}")
                                ocr_text = "(OCR 失败或未识别到文本)"
                            documents.append(Document(
                                text=f"# Slide {slide_idx} 图片 OCR\n\n{ocr_text}",
                                metadata={**base_meta, "slide": slide_idx, "shape_type": "image"}
                            ))
                            cnt_image += 1
            except Exception as e:
                logger.error(f"[PPT] Slide {slide_idx} 处理异常: {e}")
            finally:
                added = len(documents) - slide_docs_before
                logger.info(f"[PPT] Slide {slide_idx} 完成: 新增块={added}, 文本={cnt_text}, 表格={cnt_table}, 图片={cnt_image}")

    logger.info(f"[PPT] 解析完成: file={ppt_path}, 输出块数={len(documents)}")
    return documents


