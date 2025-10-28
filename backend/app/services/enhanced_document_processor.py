"""
增强型文档处理服务
支持PDF、Word、Excel、PPT的高质量解析，包括文本、表格、图片等混合内容
"""

import os
import logging
import asyncio
from typing import Dict, List, Any, Optional, Tuple
from pathlib import Path
from dataclasses import dataclass
from enum import Enum

logger = logging.getLogger(__name__)


class ChunkType(Enum):
    """文档块类型"""
    TEXT = "text"
    TABLE = "table"
    IMAGE = "image"
    TITLE = "title"
    LIST = "list"
    MIXED = "mixed"


@dataclass
class DocumentChunk:
    """文档块"""
    content: str
    chunk_type: ChunkType
    metadata: Dict[str, Any]
    has_table: bool = False
    has_image: bool = False
    page_number: Optional[int] = None
    section: Optional[str] = None


@dataclass
class ProcessedDocument:
    """处理后的文档"""
    file_path: str
    file_type: str
    chunks: List[DocumentChunk]
    metadata: Dict[str, Any]
    quality_score: float = 0.0
    processing_time: float = 0.0


class ProcessingProgress:
    """处理进度追踪"""
    
    def __init__(self):
        self.total_steps = 0
        self.current_step = 0
        self.status = "pending"
        self.logs = []
        self.start_time = None
        self.end_time = None
    
    def update(self, step: str, progress: float):
        """更新处理进度"""
        self.current_step += 1
        self.logs.append(f"步骤 {self.current_step}: {step}")
        logger.info(f"文档处理进度: {step} ({progress:.1%})")
    
    def set_status(self, status: str):
        """设置状态"""
        self.status = status
        logger.info(f"处理状态: {status}")


class EnhancedDocumentProcessor:
    """增强型文档处理器"""
    
    def __init__(self):
        self.ocr_enabled = True
        self.extract_images = True
        self.extract_tables = True
        self.progress = ProcessingProgress()
        
        # 初始化OCR服务
        self._init_ocr()
    
    def _init_ocr(self):
        """初始化OCR服务"""
        try:
            from .enhanced_ocr_service import EnhancedOCRService
            self.ocr_service = EnhancedOCRService()
            
            if self.ocr_service.is_available():
                logger.info("✅ 增强OCR服务初始化成功")
            else:
                logger.warning("增强OCR服务不可用")
                self.ocr_enabled = False
                
        except Exception as e:
            logger.warning(f"增强OCR服务初始化失败: {e}")
            # 回退到基础OCR
            try:
                import pytesseract
                pytesseract.get_tesseract_version()
                logger.info("✅ 基础OCR服务初始化成功")
            except Exception as e2:
                logger.warning(f"基础OCR服务也失败: {e2}")
                self.ocr_enabled = False
    
    async def process_document(self, file_path: str) -> ProcessedDocument:
        """
        处理文档的主入口
        
        Args:
            file_path: 文件路径
            
        Returns:
            ProcessedDocument: 处理后的文档
        """
        import time
        start_time = time.time()
        
        try:
            self.progress.set_status("processing")
            self.progress.start_time = start_time
            
            # 1. 检测文件类型
            file_type = self._detect_file_type(file_path)
            self.progress.update("检测文件类型", 0.1)
            
            # 2. 选择处理策略
            if file_type == 'pdf':
                result = await self._process_pdf_advanced(file_path)
            elif file_type in ['docx', 'doc']:
                result = await self._process_word_advanced(file_path)
            elif file_type in ['xlsx', 'xls']:
                result = await self._process_excel_advanced(file_path)
            elif file_type in ['pptx', 'ppt']:
                result = await self._process_ppt_advanced(file_path)
            elif file_type in ['txt', 'md']:
                result = await self._process_text_advanced(file_path)
            else:
                raise ValueError(f"不支持的文件类型: {file_type}")
            
            # 3. 智能分块
            self.progress.update("智能分块", 0.8)
            logger.info(f"准备分块，elements数量: {len(result.get('elements', []))}")
            chunks = self._smart_chunking(result['elements'], file_type)
            logger.info(f"分块完成，chunks数量: {len(chunks)}")
            
            # 4. 质量检查
            self.progress.update("质量检查", 0.9)
            quality_score = self._calculate_quality_score(chunks)
            logger.info(f"质量分数: {quality_score:.2f}, chunks数量: {len(chunks)}")
            
            processing_time = time.time() - start_time
            
            processed_doc = ProcessedDocument(
                file_path=file_path,
                file_type=file_type,
                chunks=chunks,
                metadata=result['metadata'],
                quality_score=quality_score,
                processing_time=processing_time
            )
            
            self.progress.set_status("completed")
            self.progress.end_time = time.time()
            
            logger.info(f"文档处理完成: {file_path}, 耗时: {processing_time:.2f}s, 质量分数: {quality_score:.2f}")
            return processed_doc
            
        except Exception as e:
            self.progress.set_status("failed")
            logger.error(f"文档处理失败 {file_path}: {str(e)}")
            raise
    
    def _detect_file_type(self, file_path: str) -> str:
        """检测文件类型"""
        try:
            import magic
            mime_type = magic.from_file(file_path, mime=True)
            
            if mime_type == 'application/pdf':
                return 'pdf'
            elif mime_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 
                              'application/msword']:
                return 'docx' if file_path.endswith('.docx') else 'doc'
            elif mime_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                              'application/vnd.ms-excel']:
                return 'xlsx' if file_path.endswith('.xlsx') else 'xls'
            elif mime_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                              'application/vnd.ms-powerpoint']:
                return 'pptx' if file_path.endswith('.pptx') else 'ppt'
            elif mime_type.startswith('text/'):
                return 'txt'
            else:
                # 回退到文件扩展名
                ext = Path(file_path).suffix.lower()
                return ext[1:] if ext else 'unknown'
        except Exception as e:
            logger.warning(f"文件类型检测失败，使用扩展名: {e}")
            ext = Path(file_path).suffix.lower()
            return ext[1:] if ext else 'unknown'
    
    async def _process_pdf_advanced(self, file_path: str) -> Dict:
        """高级PDF处理"""
        try:
            from unstructured.partition.pdf import partition_pdf
            
            self.progress.update("PDF解析", 0.2)
            
            # 使用unstructured高级PDF处理，使用更稳定的策略
            elements = partition_pdf(
                filename=file_path,
                strategy="fast",             # 使用fast策略，更稳定
                infer_table_structure=True,  # 启用表格结构推断
                extract_images_in_pdf=False, # 暂时禁用图片提取，避免问题
                languages=["chi_sim", "eng"]  # 中英文OCR
            )
            
            self.progress.update("PDF内容提取", 0.4)
            
            # 提取元数据
            metadata = {
                'file_type': 'pdf',
                'file_size': os.path.getsize(file_path),
                'element_count': len(elements),
                'has_tables': any(elem.category == "Table" for elem in elements),
                'has_images': any(elem.category == "Image" for elem in elements)
            }
            
            return {
                'elements': elements,
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"PDF处理失败: {e}")
            # 回退到基础处理
            return await self._process_pdf_fallback(file_path)
    
    async def _process_pdf_fallback(self, file_path: str) -> Dict:
        """PDF处理回退方案"""
        try:
            import fitz  # PyMuPDF
            
            doc = fitz.open(file_path)
            elements = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    # 创建简单的文本元素
                    from unstructured.documents.elements import NarrativeText
                    element = NarrativeText(text.strip())
                    element.metadata = {"page_number": page_num + 1}
                    elements.append(element)
            
            doc.close()
            
            metadata = {
                'file_type': 'pdf',
                'file_size': os.path.getsize(file_path),
                'element_count': len(elements),
                'processing_method': 'fallback'
            }
            
            return {
                'elements': elements,
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"PDF回退处理也失败: {e}")
            raise
    
    async def _process_word_advanced(self, file_path: str) -> Dict:
        """高级Word文档处理"""
        try:
            from unstructured.partition.docx import partition_docx
            
            self.progress.update("Word文档解析", 0.2)
            
            elements = partition_docx(
                filename=file_path,
                infer_table_structure=True  # 启用表格结构推断
            )
            
            self.progress.update("Word内容提取", 0.4)
            
            metadata = {
                'file_type': 'word',
                'file_size': os.path.getsize(file_path),
                'element_count': len(elements),
                'has_tables': any(elem.category == "Table" for elem in elements),
                'has_images': any(elem.category == "Image" for elem in elements)
            }
            
            return {
                'elements': elements,
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Word处理失败: {e}")
            raise
    
    async def _process_excel_advanced(self, file_path: str) -> Dict:
        """高级Excel处理"""
        try:
            import pandas as pd
            from unstructured.partition.xlsx import partition_xlsx
            
            self.progress.update("Excel解析", 0.2)
            
            elements = partition_xlsx(filename=file_path)
            
            self.progress.update("Excel内容提取", 0.4)
            
            # 额外处理：读取所有工作表
            workbook = pd.ExcelFile(file_path)
            sheet_info = {}
            
            for sheet_name in workbook.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                sheet_info[sheet_name] = {
                    'rows': len(df),
                    'columns': len(df.columns),
                    'column_names': df.columns.tolist()
                }
            
            metadata = {
                'file_type': 'excel',
                'file_size': os.path.getsize(file_path),
                'element_count': len(elements),
                'sheet_count': len(workbook.sheet_names),
                'sheet_info': sheet_info
            }
            
            return {
                'elements': elements,
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Excel处理失败: {e}")
            raise
    
    async def _process_ppt_advanced(self, file_path: str) -> Dict:
        """高级PPT处理"""
        try:
            from unstructured.partition.pptx import partition_pptx
            
            self.progress.update("PPT解析", 0.2)
            
            elements = partition_pptx(filename=file_path)
            
            self.progress.update("PPT内容提取", 0.4)
            
            metadata = {
                'file_type': 'ppt',
                'file_size': os.path.getsize(file_path),
                'element_count': len(elements),
                'has_tables': any(elem.category == "Table" for elem in elements),
                'has_images': any(elem.category == "Image" for elem in elements)
            }
            
            return {
                'elements': elements,
                'metadata': metadata
            }
            
        except ImportError as e:
            logger.warning(f"PPT高级处理所需库未安装，回退到基础处理: {e}")
            # 回退到基础文本处理
            return await self._process_ppt_basic(file_path)
        except Exception as e:
            logger.error(f"PPT处理失败: {e}")
            # 回退到基础处理
            try:
                return await self._process_ppt_basic(file_path)
            except Exception as basic_error:
                logger.error(f"PPT基础处理也失败: {basic_error}")
                raise
    
    async def _process_ppt_basic(self, file_path: str) -> Dict:
        """基础PPT处理（回退方案）"""
        try:
            # 尝试使用 python-pptx 库
            try:
                from pptx import Presentation
                
                presentation = Presentation(file_path)
                content_lines = []
                
                # 提取每张幻灯片的内容
                for slide_num, slide in enumerate(presentation.slides, 1):
                    slide_text = f"幻灯片 {slide_num}:"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text += f"\n{shape.text}"
                    content_lines.append(slide_text)
                
                # 创建简单的文本元素
                from unstructured.documents.elements import NarrativeText
                element = NarrativeText("\n".join(content_lines))
                element.metadata = {
                    "page_number": 1,
                    "file_type": "ppt",
                    "file_size": os.path.getsize(file_path)
                }
                elements = [element]
                
                metadata = {
                    'file_type': 'ppt',
                    'file_size': os.path.getsize(file_path),
                    'element_count': len(elements),
                    'has_tables': False,
                    'has_images': False,
                    'slide_count': len(presentation.slides)
                }
                
                logger.info(f"使用基础PPT处理，提取了 {len(presentation.slides)} 张幻灯片")
                return {
                    'elements': elements,
                    'metadata': metadata
                }
                
            except ImportError:
                # 如果没有 python-pptx，返回空内容
                logger.warning("python-pptx 未安装，PPT文件将无法处理")
                from unstructured.documents.elements import NarrativeText
                element = NarrativeText("PPT文件无法解析：缺少必要的依赖库。请安装 python-pptx。")
                element.metadata = {
                    "file_type": "ppt",
                    "file_size": os.path.getsize(file_path),
                    "error": "missing_dependencies"
                }
                
                return {
                    'elements': [element],
                    'metadata': {
                        'file_type': 'ppt',
                        'file_size': os.path.getsize(file_path),
                        'element_count': 1,
                        'has_tables': False,
                        'has_images': False,
                        'processing_error': 'missing_dependencies'
                    }
                }
        except Exception as e:
            logger.error(f"基础PPT处理失败: {e}")
            raise
    
    async def _process_text_advanced(self, file_path: str) -> Dict:
        """高级文本处理 - 避免NLTK依赖"""
        try:
            # 使用简单的文本读取，完全避免unstructured的NLTK依赖
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            self.progress.update("文本解析", 0.2)
            
            # 直接创建简单的文本元素，不使用unstructured的复杂处理
            from unstructured.documents.elements import NarrativeText
            
            # 手动设置元素属性，避免触发NLTK
            element = NarrativeText(content)
            element.metadata = {
                "page_number": 1,
                "file_type": "text",
                "file_size": os.path.getsize(file_path)
            }
            elements = [element]
            
            self.progress.update("文本内容提取", 0.4)
            
            metadata = {
                'file_type': 'text',
                'file_size': os.path.getsize(file_path),
                'element_count': len(elements)
            }
            
            return {
                'elements': elements,
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"文本处理失败: {e}")
            raise
    
    def _smart_chunking(self, elements: List, file_type: str) -> List[DocumentChunk]:
        """智能分块策略"""
        chunks = []
        current_chunk_elements = []
        current_chunk_type = None
        
        if not elements:
            logger.warning("elements列表为空")
            return chunks
        
        logger.info(f"开始分块，elements数量: {len(elements)}")
        
        for i, element in enumerate(elements):
            element_type = self._get_element_type(element)
            
            # 尝试获取元素内容
            try:
                if hasattr(element, 'text'):
                    element_content = element.text
                elif hasattr(element, 'content'):
                    element_content = element.content
                else:
                    element_content = str(element)
            except Exception as e:
                logger.warning(f"获取element {i}内容失败: {e}")
                element_content = str(element)
            
            # 跳过空内容
            if not element_content or not element_content.strip():
                logger.debug(f"跳过空element {i}: {type(element)}")
                continue
            
            # 判断是否需要创建新块
            if self._should_create_new_chunk(element_type, current_chunk_type, current_chunk_elements):
                if current_chunk_elements:
                    # 创建当前块
                    chunk = self._create_chunk_from_elements(current_chunk_elements, current_chunk_type)
                    chunks.append(chunk)
                
                # 开始新块
                current_chunk_elements = [element]
                current_chunk_type = element_type
            else:
                current_chunk_elements.append(element)
                if current_chunk_type is None:
                    current_chunk_type = element_type
        
        # 处理最后一个块
        if current_chunk_elements:
            chunk = self._create_chunk_from_elements(current_chunk_elements, current_chunk_type)
            chunks.append(chunk)
        
        return chunks
    
    def _get_element_type(self, element) -> ChunkType:
        """获取元素类型"""
        category = getattr(element, 'category', 'NarrativeText')
        
        if category == "Table":
            return ChunkType.TABLE
        elif category == "Image":
            return ChunkType.IMAGE
        elif category == "Title":
            return ChunkType.TITLE
        elif category == "ListItem":
            return ChunkType.LIST
        else:
            return ChunkType.TEXT
    
    def _should_create_new_chunk(self, element_type: ChunkType, current_type: ChunkType, 
                                current_elements: List) -> bool:
        """判断是否应该创建新块"""
        # 表格和图片单独成块
        if element_type in [ChunkType.TABLE, ChunkType.IMAGE]:
            return True
        
        # 如果当前块太大，创建新块
        if len(current_elements) > 10:  # 增加块大小限制
            return True
        
        # 计算当前块的内容长度
        current_content_length = sum(len(str(elem)) for elem in current_elements)
        if current_content_length > 2000:  # 内容长度限制
            return True
        
        return False
    
    def _create_chunk_from_elements(self, elements: List, chunk_type: ChunkType) -> DocumentChunk:
        """从元素列表创建文档块"""
        content_parts = []
        has_table = False
        has_image = False
        page_number = None
        
        for element in elements:
            element_type = self._get_element_type(element)
            
            # 获取元素内容 - 使用text属性而不是str()
            try:
                if hasattr(element, 'text'):
                    element_content = element.text
                elif hasattr(element, 'content'):
                    element_content = element.content
                else:
                    element_content = str(element)
            except:
                element_content = str(element)
            
            # 过滤空内容
            if not element_content or not element_content.strip():
                continue
            
            if element_type == ChunkType.TABLE:
                has_table = True
                # 表格内容特殊处理
                content_parts.append(f"[表格数据]\n{element_content}")
            elif element_type == ChunkType.IMAGE:
                has_image = True
                content_parts.append(f"[图片内容]\n{element_content}")
            else:
                content_parts.append(element_content)
            
            # 获取页码
            if hasattr(element, 'metadata') and element.metadata:
                try:
                    if hasattr(element.metadata, 'get'):
                        page_number = element.metadata.get('page_number')
                    else:
                        # 如果是ElementMetadata对象，直接访问属性
                        page_number = getattr(element.metadata, 'page_number', None)
                except:
                    page_number = None
        
        content = '\n\n'.join(content_parts)
        
        metadata = {
            'element_count': len(elements),
            'chunk_type': chunk_type.value,
            'created_at': str(asyncio.get_event_loop().time())
        }
        
        return DocumentChunk(
            content=content,
            chunk_type=chunk_type,
            metadata=metadata,
            has_table=has_table,
            has_image=has_image,
            page_number=page_number
        )
    
    def _calculate_quality_score(self, chunks: List[DocumentChunk]) -> float:
        """计算文档质量分数"""
        if not chunks:
            return 0.0
        
        total_score = 0.0
        
        for chunk in chunks:
            chunk_score = 0.0
            
            # 内容长度评分
            if len(chunk.content) > 50:
                chunk_score += 0.3
            
            # 内容完整性评分
            if chunk.content.strip():
                chunk_score += 0.3
            
            # 表格处理评分
            if chunk.has_table:
                chunk_score += 0.2
            
            # 图片处理评分
            if chunk.has_image:
                chunk_score += 0.2
            
            total_score += chunk_score
        
        return total_score / len(chunks)
    
    def get_processing_logs(self) -> List[str]:
        """获取处理日志"""
        return self.progress.logs
    
    def get_processing_status(self) -> Dict[str, Any]:
        """获取处理状态"""
        return {
            'status': self.progress.status,
            'current_step': self.progress.current_step,
            'total_steps': self.progress.total_steps,
            'logs': self.progress.logs,
            'start_time': self.progress.start_time,
            'end_time': self.progress.end_time
        }
