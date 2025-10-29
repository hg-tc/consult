"""
PPT 处理器
提取幻灯片文本与图片，使用 Tesseract OCR 对图片进行识别，合并为 Markdown
"""

import os
from pathlib import Path
from typing import Dict, Any, List
import logging

from pptx import Presentation

from app.services.ocr_service import OCRService

logger = logging.getLogger(__name__)


def process_ppt_to_markdown(file_path: str) -> Dict[str, Any]:
    if not os.path.exists(file_path):
        raise FileNotFoundError(file_path)

    prs = Presentation(file_path)
    ocr = OCRService()

    slides_result: List[Dict[str, Any]] = []
    md_parts: List[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        images_texts: List[str] = []

        # 文本
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                txt = shape.text.strip()
                if txt:
                    texts.append(txt)

        # 图片 OCR（导出为临时文件再处理）
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    ext = image.ext
                    tmp_path = Path(file_path).with_suffix("")
                    out_path = f"{tmp_path}_slide{idx}_img_{image.sha1}.{ext}"
                    with open(out_path, "wb") as f:
                        f.write(image.blob)
                    try:
                        text = ocr.extract_text_from_image(out_path)
                        if isinstance(text, dict):
                            t = text.get("text", "").strip()
                        else:
                            t = str(text).strip()
                        if t:
                            images_texts.append(t)
                    finally:
                        try:
                            os.remove(out_path)
                        except Exception:
                            pass
                except Exception as e:
                    logger.warning(f"PPT 图片OCR失败: {e}")

        # 合并 Markdown
        slide_md_parts: List[str] = [f"## 幻灯片 {idx}"]
        if texts:
            slide_md_parts.append("\n".join(texts))
        if images_texts:
            slide_md_parts.append("### 图片OCR\n" + "\n\n".join(images_texts))

        slide_md = "\n\n".join(slide_md_parts)
        md_parts.append(slide_md)

        slides_result.append({
            "slide_number": idx,
            "text_content": "\n".join(texts),
            "ocr_content": "\n\n".join(images_texts),
            "combined_content": slide_md,
        })

    final_md = "\n\n".join(md_parts)
    return {
        "file_type": "powerpoint",
        "slides": slides_result,
        "content": final_md,
        "metadata": {
            "slide_count": len(slides_result),
            "file_size": os.path.getsize(file_path)
        }
    }


