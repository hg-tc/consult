"""
文件处理服务
支持多种文件格式的解析和内容提取
"""

import os
import logging
from typing import Dict, List, Any, Optional
from pathlib import Path

# 文件处理库
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
import pandas as pd

logger = logging.getLogger(__name__)


class FileProcessor:
    """文件处理器"""

    def __init__(self, upload_dir: str = "uploads"):
        self.upload_dir = Path(upload_dir)
        self.upload_dir.mkdir(exist_ok=True)

    def process_file(self, file_path: str) -> Dict[str, Any]:
        """
        处理上传的文件，返回解析后的内容和元数据

        Args:
            file_path: 文件路径

        Returns:
            Dict包含：
            - content: 提取的文本内容
            - metadata: 文件元数据
            - chunks: 文本分块（如果适用）
            - file_type: 文件类型
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        file_ext = Path(file_path).suffix.lower()

        try:
            if file_ext == '.pdf':
                # PDF -> Markdown（优先使用 marker-pdf）
                try:
                    from app.services.pdf_enhanced_processor import process_pdf_to_markdown
                    return process_pdf_to_markdown(file_path)
                except Exception as e:
                    logger.warning(f"PDF增强处理器失败，使用基础处理: {e}")
                    return self._process_pdf(file_path)
            elif file_ext in ['.docx', '.doc']:
                # WORD -> Markdown（文本+表格+图片OCR）
                if file_ext == '.docx':
                    try:
                        from app.services.word_enhanced_processor import process_word_to_markdown
                        return process_word_to_markdown(file_path)
                    except Exception as e:
                        logger.warning(f"WORD增强处理器失败，使用基础处理: {e}")
                        return self._process_word(file_path)
                else:
                    # 旧版.doc文件使用基础处理
                    return self._process_word(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                # Excel -> Markdown
                try:
                    from app.services.excel_parser import parse_excel_to_markdown
                    return parse_excel_to_markdown(file_path)
                except Exception:
                    return self._process_excel(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                # PPT -> Markdown（含图片OCR）
                try:
                    from app.services.ppt_processor import process_ppt_to_markdown
                    return process_ppt_to_markdown(file_path)
                except Exception:
                    return self._process_powerpoint(file_path)
            elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
                return self._process_image(file_path)
            elif file_ext in ['.txt', '.md']:
                return self._process_text(file_path)
            else:
                raise ValueError(f"不支持的文件类型: {file_ext}")

        except Exception as e:
            logger.error(f"文件处理失败 {file_path}: {str(e)}")
            raise

    def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """处理PDF文件"""
        content = []
        metadata = {}

        try:
            doc = fitz.open(file_path)

            # 提取元数据
            metadata = doc.metadata
            metadata['page_count'] = len(doc)
            metadata['file_size'] = os.path.getsize(file_path)

            # 提取文本内容
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                if text.strip():
                    content.append({
                        'page': page_num + 1,
                        'content': text.strip()
                    })

            doc.close()

        except Exception as e:
            logger.error(f"PDF处理失败: {str(e)}")
            raise

        return {
            'content': '\n\n'.join([item['content'] for item in content]),
            'metadata': metadata,
            'chunks': content,
            'file_type': 'pdf'
        }

    def _process_word(self, file_path: str) -> Dict[str, Any]:
        """处理Word文档"""
        content = []
        metadata = {}

        try:
            doc = DocxDocument(file_path)

            # 提取元数据
            metadata['paragraph_count'] = len(doc.paragraphs)
            metadata['file_size'] = os.path.getsize(file_path)

            # 提取文本内容
            for para in doc.paragraphs:
                if para.text.strip():
                    content.append(para.text.strip())

            # 提取表格内容（如果有）
            tables_content = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    tables_content.append(table_data)

            if tables_content:
                content.extend(['表格数据:'] + [str(table) for table in tables_content])

        except Exception as e:
            logger.error(f"Word处理失败: {str(e)}")
            raise

        return {
            'content': '\n\n'.join(content),
            'metadata': metadata,
            'chunks': [{'content': para, 'type': 'paragraph'} for para in content],
            'file_type': 'word'
        }

    def _process_excel(self, file_path: str) -> Dict[str, Any]:
        """处理Excel文件"""
        content = []
        metadata = {}
        sheets_data = {}

        try:
            workbook = load_workbook(file_path, read_only=True)

            # 提取元数据
            metadata['sheet_count'] = len(workbook.sheetnames)
            metadata['file_size'] = os.path.getsize(file_path)

            # 处理每个工作表
            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]

                # 转换为DataFrame
                data = []
                for row in worksheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):  # 跳过空行
                        data.append(list(row))

                if data:
                    df = pd.DataFrame(data[1:], columns=data[0] if len(data) > 1 else None)
                    sheets_data[sheet_name] = df.to_dict('records')

                    # 生成文本描述
                    content.append(f"工作表: {sheet_name}")
                    content.append(df.to_string())
                    content.append("")

            workbook.close()

        except Exception as e:
            logger.error(f"Excel处理失败: {str(e)}")
            raise

        return {
            'content': '\n'.join(content),
            'metadata': metadata,
            'chunks': [
                {'sheet': sheet, 'data': data}
                for sheet, data in sheets_data.items()
            ],
            'file_type': 'excel',
            'sheets_data': sheets_data
        }

    def _process_powerpoint(self, file_path: str) -> Dict[str, Any]:
        """处理PowerPoint文件"""
        content = []
        metadata = {}

        try:
            presentation = Presentation(file_path)

            # 提取元数据
            metadata['slide_count'] = len(presentation.slides)
            metadata['file_size'] = os.path.getsize(file_path)

            # 提取幻灯片内容
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_content = []

                # 提取标题
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content.append(shape.text.strip())

                if slide_content:
                    content.append(f"幻灯片 {slide_num}:")
                    content.append('\n'.join(slide_content))
                    content.append("")

        except Exception as e:
            logger.error(f"PowerPoint处理失败: {str(e)}")
            raise

        return {
            'content': '\n'.join(content),
            'metadata': metadata,
            'chunks': [
                {'slide': i+1, 'content': slide_text}
                for i, slide_text in enumerate(content)
                if slide_text.startswith('幻灯片')
            ],
            'file_type': 'powerpoint'
        }

    def _process_text(self, file_path: str) -> Dict[str, Any]:
        """处理文本文件"""
        metadata = {}

        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            metadata['file_size'] = os.path.getsize(file_path)
            metadata['encoding'] = 'utf-8'

            # 简单分段
            paragraphs = [para.strip() for para in content.split('\n\n') if para.strip()]

        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                with open(file_path, 'r', encoding='gbk') as f:
                    content = f.read()
                metadata['encoding'] = 'gbk'
                paragraphs = [para.strip() for para in content.split('\n\n') if para.strip()]
            except Exception as e:
                raise ValueError(f"无法解析文本文件编码: {str(e)}")

        return {
            'content': content,
            'metadata': metadata,
            'chunks': [{'content': para, 'type': 'paragraph'} for para in paragraphs],
            'file_type': 'text'
        }

    def _process_image(self, file_path: str) -> Dict[str, Any]:
        """处理图片文件，使用本地 Tesseract OCR 输出 Markdown。"""
        try:
            import time
            start_time = time.time()
            
            logger.info(f"📷 开始处理图片文件: {file_path}")
            
            from app.services.ocr_service import OCRService
            ocr = OCRService()
            
            if not ocr.is_available():
                logger.warning(f"⚠️ OCR服务不可用，无法处理图片: {file_path}")
                return {
                    'content': "### 图片OCR\n\n⚠️ OCR服务不可用，无法提取文字",
                    'metadata': {
                        'file_size': os.path.getsize(file_path),
                        'image_path': file_path,
                        'ocr_available': False
                    },
                    'chunks': [],
                    'file_type': 'image'
                }
            
            # 调用同步版本的 OCR 方法
            logger.info(f"🔄 调用OCR服务提取文字: {file_path}")
            text = ocr.extract_text_from_image(file_path)
            
            elapsed_time = time.time() - start_time
            logger.info(f"⏱️ 图片OCR处理耗时: {elapsed_time:.2f}秒, 提取文字长度: {len(text)}")
            
            # 构建 Markdown 内容
            if text and text.strip():
                md = f"### 图片OCR识别结果\n\n{text.strip()}"
                logger.info(f"✅ 图片OCR成功: {file_path}, 文字预览: {text[:100]}...")
            else:
                md = "### 图片OCR识别结果\n\n⚠️ 未识别到任何文字（可能图片中没有文字或图片质量不佳）"
                logger.warning(f"⚠️ 图片OCR未提取到文字: {file_path}")
            
            return {
                'content': md,
                'metadata': {
                    'file_size': os.path.getsize(file_path),
                    'image_path': file_path,
                    'ocr_available': True,
                    'ocr_text_length': len(text),
                    'processing_time': elapsed_time
                },
                'chunks': [{'content': text.strip(), 'type': 'ocr'}] if text else [],
                'file_type': 'image'
            }
        except Exception as e:
            logger.error(f"❌ 图片处理失败: {file_path}, 错误: {e}", exc_info=True)
            raise

    def get_supported_formats(self) -> List[str]:
        """获取支持的文件格式"""
        return ['.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.txt', '.md',
                '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']
