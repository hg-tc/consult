"""
文档生成服务
基于模板和数据生成各种格式的文档
"""

import os
import json
import logging
from typing import Dict, List, Any, Optional, BinaryIO
from pathlib import Path
from io import BytesIO

from docx import Document as DocxDocument
from docx.shared import Inches
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)


class DocumentGenerator:
    """文档生成器"""

    def __init__(self, template_dir: str = "templates", output_dir: str = "generated"):
        self.template_dir = Path(template_dir)
        self.output_dir = Path(output_dir)
        self.template_dir.mkdir(exist_ok=True)
        self.output_dir.mkdir(exist_ok=True)

    def generate_from_template(
        self,
        template_path: str,
        data: Dict[str, Any],
        output_format: str = "docx"
    ) -> str:
        """
        基于模板生成文档

        Args:
            template_path: 模板文件路径
            data: 填充数据
            output_format: 输出格式 (docx, pdf, pptx)

        Returns:
            生成的文档文件路径
        """
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"模板文件不存在: {template_path}")

        template_ext = Path(template_path).suffix.lower()

        if template_ext == '.docx':
            return self._generate_from_docx_template(template_path, data, output_format)
        elif template_ext in ['.pptx', '.ppt']:
            return self._generate_from_pptx_template(template_path, data, output_format)
        else:
            raise ValueError(f"不支持的模板格式: {template_ext}")

    def _generate_from_docx_template(
        self,
        template_path: str,
        data: Dict[str, Any],
        output_format: str
    ) -> str:
        """基于Word模板生成文档"""
        try:
            # 读取模板
            template_doc = DocxDocument(template_path)

            # 替换模板中的变量
            self._replace_docx_variables(template_doc, data)

            # 生成文件名
            timestamp = str(int(__import__('time').time()))
            filename = f"generated_{timestamp}.docx"
            output_path = self.output_dir / filename

            # 保存文档
            template_doc.save(str(output_path))

            # 如果需要转换为其他格式
            if output_format == 'pdf':
                return self._convert_docx_to_pdf(str(output_path))
            elif output_format == 'pptx':
                return self._convert_docx_to_pptx(str(output_path))

            return str(output_path)

        except Exception as e:
            logger.error(f"Word文档生成失败: {str(e)}")
            raise

    def _replace_docx_variables(self, doc: DocxDocument, data: Dict[str, Any]):
        """替换Word文档中的变量"""
        # 替换段落中的文本
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                for key, value in data.items():
                    placeholder = f"{{{key}}}"
                    if placeholder in run.text:
                        run.text = run.text.replace(placeholder, str(value))

        # 替换表格中的内容
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            for key, value in data.items():
                                placeholder = f"{{{key}}}"
                                if placeholder in run.text:
                                    run.text = run.text.replace(placeholder, str(value))

    def _generate_from_pptx_template(
        self,
        template_path: str,
        data: Dict[str, Any],
        output_format: str
    ) -> str:
        """基于PPT模板生成演示文稿"""
        try:
            # 读取模板
            presentation = Presentation(template_path)

            # 替换幻灯片内容
            self._replace_pptx_variables(presentation, data)

            # 生成文件名
            timestamp = str(int(__import__('time').time()))
            filename = f"generated_{timestamp}.pptx"
            output_path = self.output_dir / filename

            # 保存演示文稿
            presentation.save(str(output_path))

            return str(output_path)

        except Exception as e:
            logger.error(f"PPT文档生成失败: {str(e)}")
            raise

    def _replace_pptx_variables(self, presentation: Presentation, data: Dict[str, Any]):
        """替换PPT中的变量"""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    for key, value in data.items():
                        placeholder = f"{{{key}}}"
                        if placeholder in shape.text:
                            shape.text = shape.text.replace(placeholder, str(value))

    def _convert_docx_to_pdf(self, docx_path: str) -> str:
        """将Word文档转换为PDF"""
        try:
            # 这里需要使用第三方库如docx2pdf或类似的工具
            # 暂时返回原路径，表示转换失败
            logger.warning("Word转PDF功能需要额外安装docx2pdf库")
            return docx_path

        except Exception as e:
            logger.error(f"Word转PDF失败: {str(e)}")
            return docx_path

    def _convert_docx_to_pptx(self, docx_path: str) -> str:
        """将Word文档转换为PPT"""
        try:
            # 这里需要复杂的文档结构分析和转换逻辑
            # 暂时返回原路径，表示转换失败
            logger.warning("Word转PPT功能需要复杂的文档结构分析")
            return docx_path

        except Exception as e:
            logger.error(f"Word转PPT失败: {str(e)}")
            return docx_path

    def create_document_from_data(
        self,
        data: Dict[str, Any],
        template_type: str = "report",
        output_format: str = "docx"
    ) -> str:
        """
        根据数据创建文档（不基于模板）

        Args:
            data: 文档数据
            template_type: 模板类型 (report, contract, presentation等)
            output_format: 输出格式

        Returns:
            生成的文档文件路径
        """
        if template_type == "report":
            return self._create_report_document(data, output_format)
        elif template_type == "contract":
            return self._create_contract_document(data, output_format)
        else:
            raise ValueError(f"不支持的模板类型: {template_type}")

    def _create_report_document(self, data: Dict[str, Any], output_format: str) -> str:
        """创建报告文档"""
        try:
            if output_format == "docx":
                return self._create_docx_report(data)
            elif output_format == "pdf":
                return self._create_pdf_report(data)
            else:
                raise ValueError(f"不支持的输出格式: {output_format}")

        except Exception as e:
            logger.error(f"报告文档生成失败: {str(e)}")
            raise

    def _create_docx_report(self, data: Dict[str, Any]) -> str:
        """创建Word格式的报告"""
        doc = DocxDocument()

        # 添加标题
        title = data.get('title', '报告标题')
        doc.add_heading(title, 0)

        # 添加基本信息
        if 'summary' in data:
            doc.add_paragraph(data['summary'])

        # 添加章节
        if 'sections' in data:
            for section in data['sections']:
                doc.add_heading(section.get('title', '章节标题'), level=1)
                if 'content' in section:
                    doc.add_paragraph(section['content'])

        # 生成文件名
        timestamp = str(int(__import__('time').time()))
        filename = f"report_{timestamp}.docx"
        output_path = self.output_dir / filename

        doc.save(str(output_path))
        return str(output_path)

    def _create_pdf_report(self, data: Dict[str, Any]) -> str:
        """创建PDF格式的报告"""
        try:
            # 暂时跳过PDF生成，返回Word文档路径
            logger.warning("PDF生成功能暂时不可用，使用Word文档替代")

            # 生成Word文档作为替代
            return self._create_docx_report(data)

        except Exception as e:
            logger.error(f"PDF报告生成失败: {str(e)}")
            raise

    def _create_contract_document(self, data: Dict[str, Any], output_format: str) -> str:
        """创建合同文档"""
        # 这里可以实现合同模板的生成逻辑
        # 暂时返回报告生成的文档
        return self._create_report_document(data, output_format)

    def get_supported_formats(self) -> List[str]:
        """获取支持的输出格式"""
        return ['docx', 'pdf', 'pptx']

    def get_template_variables(self, template_path: str) -> List[str]:
        """提取模板中的变量"""
        try:
            if template_path.endswith('.docx'):
                doc = DocxDocument(template_path)
                variables = set()

                # 从段落中提取变量
                for paragraph in doc.paragraphs:
                    for run in paragraph.runs:
                        # 简单的变量提取逻辑
                        import re
                        matches = re.findall(r'\{([^}]+)\}', run.text)
                        variables.update(matches)

                return list(variables)

            else:
                return []

        except Exception as e:
            logger.error(f"模板变量提取失败: {str(e)}")
            return []

    def validate_template_data(self, template_path: str, data: Dict[str, Any]) -> Dict[str, Any]:
        """验证模板数据完整性"""
        required_vars = self.get_template_variables(template_path)
        missing_vars = []
        extra_vars = []

        # 检查缺少的变量
        for var in required_vars:
            if var not in data:
                missing_vars.append(var)

        # 检查额外的变量
        for key in data.keys():
            if key not in required_vars:
                extra_vars.append(key)

        return {
            'valid': len(missing_vars) == 0,
            'missing_variables': missing_vars,
            'extra_variables': extra_vars,
            'required_variables': required_vars
        }
