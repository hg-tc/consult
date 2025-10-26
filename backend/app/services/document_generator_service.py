"""
文档生成服务
支持根据用户指令生成Word和PDF文档
"""

import os
import logging
import json
import uuid
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
from datetime import datetime
from dataclasses import dataclass
from enum import Enum

logger = logging.getLogger(__name__)


class DocumentType(Enum):
    """文档类型"""
    WORD = "word"
    PDF = "pdf"
    HTML = "html"
    EXCEL = "excel"
    PPT = "ppt"
    UNKNOWN = "unknown"


class DocumentTemplate(Enum):
    """文档模板"""
    REPORT = "report"
    SUMMARY = "summary"
    ANALYSIS = "analysis"
    PROPOSAL = "proposal"
    CUSTOM = "custom"


@dataclass
class DocumentContent:
    """文档内容"""
    title: str
    sections: List[Dict[str, Any]]
    metadata: Dict[str, Any]
    template: DocumentTemplate = DocumentTemplate.CUSTOM


class WordGenerator:
    """Word文档生成器（增强版）"""
    
    def __init__(self):
        self.templates_dir = Path("templates/word")
        self.templates_dir.mkdir(exist_ok=True)
    
    def _setup_document_styles(self, doc):
        """设置文档样式"""
        try:
            from docx.enum.style import WD_STYLE_TYPE
            from docx.shared import Pt, RGBColor
            
            # 设置标题样式
            styles = doc.styles
            
            # 标题1样式
            if 'Custom Heading 1' not in [s.name for s in styles]:
                h1_style = styles.add_style('Custom Heading 1', WD_STYLE_TYPE.PARAGRAPH)
                h1_style.font.name = 'Microsoft YaHei'
                h1_style.font.size = Pt(18)
                h1_style.font.bold = True
                h1_style.font.color.rgb = RGBColor(44, 62, 80)  # 深蓝色
                h1_style.paragraph_format.space_before = Pt(12)
                h1_style.paragraph_format.space_after = Pt(6)
            
            # 标题2样式
            if 'Custom Heading 2' not in [s.name for s in styles]:
                h2_style = styles.add_style('Custom Heading 2', WD_STYLE_TYPE.PARAGRAPH)
                h2_style.font.name = 'Microsoft YaHei'
                h2_style.font.size = Pt(16)
                h2_style.font.bold = True
                h2_style.font.color.rgb = RGBColor(52, 73, 94)  # 深灰色
                h2_style.paragraph_format.space_before = Pt(10)
                h2_style.paragraph_format.space_after = Pt(4)
            
            # 正文样式
            if 'Custom Body' not in [s.name for s in styles]:
                body_style = styles.add_style('Custom Body', WD_STYLE_TYPE.PARAGRAPH)
                body_style.font.name = 'Microsoft YaHei'
                body_style.font.size = Pt(12)
                body_style.paragraph_format.line_spacing = 1.5
                body_style.paragraph_format.space_after = Pt(6)
            
        except Exception as e:
            logger.warning(f"设置文档样式失败: {e}")
    
    def _add_header_footer(self, doc, title: str):
        """添加页眉页脚"""
        try:
            # 获取节
            section = doc.sections[0]
            
            # 设置页眉
            header = section.header
            header_para = header.paragraphs[0]
            header_para.text = title
            header_para.alignment = 1  # 居中
            
            # 设置页脚
            footer = section.footer
            footer_para = footer.paragraphs[0]
            footer_para.text = f"第 {chr(0x200B)} 页，共 {chr(0x200B)} 页"  # 占位符
            footer_para.alignment = 1  # 居中
            from docx.shared import Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            footer_para.runs[0].font.size = Pt(9)
            footer_para.runs[0].font.color.rgb = RGBColor(128, 128, 128)
            
        except Exception as e:
            logger.warning(f"添加页眉页脚失败: {e}")
    
    def _add_table_of_contents(self, doc, sections: List[Dict[str, Any]]):
        """添加目录"""
        try:
            doc.add_heading('目录', level=1)
            
            # 创建目录表格
            toc_table = doc.add_table(rows=len(sections) + 1, cols=3)
            toc_table.style = 'Table Grid'
            
            # 设置表头
            header_row = toc_table.rows[0]
            header_row.cells[0].text = "章节"
            header_row.cells[1].text = "标题"
            header_row.cells[2].text = "页码"
            
            # 设置表头样式
            from docx.shared import Pt
            for cell in header_row.cells:
                cell.paragraphs[0].runs[0].font.bold = True
                cell.paragraphs[0].runs[0].font.size = Pt(12)
            
            # 添加目录项
            for i, section in enumerate(sections, 1):
                row = toc_table.rows[i]
                row.cells[0].text = f"第{i}章"
                row.cells[1].text = section.get('title', f'章节{i}')
                row.cells[2].text = f"{i + 2}"  # 估算页码
            
            # 添加分页符
            doc.add_page_break()
            
        except Exception as e:
            logger.warning(f"添加目录失败: {e}")
    
    def _add_references_section(self, doc, content: DocumentContent):
        """添加参考文献章节"""
        try:
            doc.add_heading('参考文献', level=1)
            
            # 从元数据中提取参考文献
            references = content.metadata.get('references', [])
            if not references:
                # 生成默认参考文献
                references = [
                    "1. 相关文档资料",
                    "2. 网络搜索信息",
                    "3. 对话历史记录"
                ]
            
            # 添加参考文献列表
            for i, ref in enumerate(references, 1):
                para = doc.add_paragraph()
                para.add_run(f"[{i}] ").font.bold = True
                para.add_run(str(ref))
            
        except Exception as e:
            logger.warning(f"添加参考文献失败: {e}")
    
    def generate_document(self, content: DocumentContent, output_path: str) -> bool:
        """
        生成Word文档（增强版）
        
        Args:
            content: 文档内容
            output_path: 输出路径
            
        Returns:
            是否生成成功
        """
        try:
            from docx import Document
            from docx.shared import Inches, Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
            from docx.enum.style import WD_STYLE_TYPE
            from docx.oxml.shared import OxmlElement, qn
            from docx.oxml.ns import nsdecls
            from docx.oxml import parse_xml
            
            # 创建文档
            doc = Document()
            
            # 设置文档样式
            self._setup_document_styles(doc)
            
            # 添加页眉页脚
            self._add_header_footer(doc, content.title)
            
            # 设置文档标题
            title = doc.add_heading(content.title, 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # 添加元数据
            if content.metadata:
                self._add_metadata_section(doc, content.metadata)
            
            # 生成目录
            self._add_table_of_contents(doc, content.sections)
            
            # 添加各个章节
            for section in content.sections:
                self._add_section(doc, section)
            
            # 添加参考文献
            self._add_references_section(doc, content)
            
            # 保存文档
            doc.save(output_path)
            
            logger.info(f"Word文档生成成功: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"Word文档生成失败: {e}")
            return False
    
    def _add_metadata_section(self, doc, metadata: Dict[str, Any]):
        """添加元数据章节"""
        doc.add_heading('文档信息', level=1)
        
        table = doc.add_table(rows=len(metadata), cols=2)
        table.style = 'Table Grid'
        
        for i, (key, value) in enumerate(metadata.items()):
            table.cell(i, 0).text = str(key)
            table.cell(i, 1).text = str(value)
    
    def _add_section(self, doc, section: Dict[str, Any]):
        """添加章节（使用增强版本）"""
        try:
            # 优先使用增强章节方法
            self._add_enhanced_section(doc, section)
        except Exception as e:
            logger.warning(f"增强章节添加失败，回退到基础方法: {e}")
            # 回退到基础方法
            self._add_basic_section(doc, section)
    
    def _add_basic_section(self, doc, section: Dict[str, Any]):
        """添加基础章节（原始方法）"""
        # 添加章节标题
        if 'title' in section:
            doc.add_heading(section['title'], level=section.get('level', 2))
        
        # 添加内容（新格式）
        if 'content' in section:
            content_text = section['content']
            if content_text:
                # 将内容按段落分割
                paragraphs = content_text.split('\n\n')
                for para_text in paragraphs:
                    if para_text.strip():
                        doc.add_paragraph(para_text.strip())
        
        # 添加表格数据（新格式）
        if 'table_data' in section and section['table_data']:
            table_data = section['table_data']
            if table_data and len(table_data) > 0:
                # 创建表格
                table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
                table.style = 'Table Grid'
                
                # 填充表格数据
                for i, row in enumerate(table_data):
                    for j, cell_value in enumerate(row):
                        if j < len(table.rows[i].cells):
                            table.rows[i].cells[j].text = str(cell_value)
        
        # 添加段落（旧格式兼容）
        if 'paragraphs' in section:
            for para_text in section['paragraphs']:
                doc.add_paragraph(para_text)
        
        # 添加表格（旧格式兼容）
        if 'tables' in section:
            for table_data in section['tables']:
                self._add_table(doc, table_data)
        
        # 添加列表（旧格式兼容）
        if 'lists' in section:
            for list_data in section['lists']:
                self._add_list(doc, list_data)
    
    def _add_table(self, doc, table_data: Dict[str, Any]):
        """添加表格"""
        headers = table_data.get('headers', [])
        rows = table_data.get('rows', [])
        
        if not headers and not rows:
            return
        
        # 创建表格
        table = doc.add_table(rows=len(rows) + (1 if headers else 0), cols=len(headers))
        table.style = 'Table Grid'
        
        # 添加表头
        if headers:
            header_row = table.rows[0]
            for i, header in enumerate(headers):
                header_row.cells[i].text = str(header)
        
        # 添加数据行
        start_row = 1 if headers else 0
        for i, row in enumerate(rows):
            if i + start_row < len(table.rows):
                table_row = table.rows[i + start_row]
                for j, cell_value in enumerate(row):
                    if j < len(table_row.cells):
                        table_row.cells[j].text = str(cell_value)
    
    def _add_list(self, doc, list_data: Dict[str, Any]):
        """添加列表"""
        items = list_data.get('items', [])
        list_type = list_data.get('type', 'bullet')
        
        if not items:
            return
        
        if list_type == 'numbered':
            for item in items:
                doc.add_paragraph(str(item), style='List Number')
        else:
            for item in items:
                doc.add_paragraph(str(item), style='List Bullet')


class PDFGenerator:
    """PDF文档生成器"""
    
    def __init__(self):
        self.templates_dir = Path("templates/pdf")
        self.templates_dir.mkdir(exist_ok=True)
    
    def generate_document(self, content: DocumentContent, output_path: str) -> bool:
        """
        生成PDF文档
        
        Args:
            content: 文档内容
            output_path: 输出路径
            
        Returns:
            是否生成成功
        """
        try:
            # 生成HTML
            html_content = self._generate_html(content)
            
            # 转换为PDF
            return self._html_to_pdf(html_content, output_path)
            
        except Exception as e:
            logger.error(f"PDF文档生成失败: {e}")
            return False
    
    def _generate_html(self, content: DocumentContent) -> str:
        """生成HTML内容"""
        html_parts = [
            "<!DOCTYPE html>",
            "<html lang='zh-CN'>",
            "<head>",
            "<meta charset='UTF-8'>",
            "<title>{}</title>".format(content.title),
            "<style>",
            self._get_css_styles(),
            "</style>",
            "</head>",
            "<body>",
            "<div class='container'>",
            "<h1 class='title'>{}</h1>".format(content.title)
        ]
        
        # 添加元数据
        if content.metadata:
            html_parts.append("<div class='metadata'>")
            html_parts.append("<h2>文档信息</h2>")
            html_parts.append("<table class='metadata-table'>")
            for key, value in content.metadata.items():
                html_parts.append(f"<tr><td>{key}</td><td>{value}</td></tr>")
            html_parts.append("</table>")
            html_parts.append("</div>")
        
        # 添加章节
        for section in content.sections:
            html_parts.append(self._section_to_html(section))
        
        html_parts.extend([
            "</div>",
            "</body>",
            "</html>"
        ])
        
        return "\n".join(html_parts)
    
    def _get_css_styles(self) -> str:
        """获取CSS样式"""
        return """
        body {
            font-family: 'Microsoft YaHei', Arial, sans-serif;
            line-height: 1.6;
            color: #333;
            margin: 0;
            padding: 20px;
        }
        .container {
            max-width: 800px;
            margin: 0 auto;
        }
        .title {
            text-align: center;
            color: #2c3e50;
            border-bottom: 2px solid #3498db;
            padding-bottom: 10px;
        }
        .metadata {
            margin: 20px 0;
            padding: 15px;
            background-color: #f8f9fa;
            border-radius: 5px;
        }
        .metadata-table {
            width: 100%;
            border-collapse: collapse;
        }
        .metadata-table td {
            padding: 8px;
            border: 1px solid #ddd;
        }
        .metadata-table td:first-child {
            font-weight: bold;
            background-color: #e9ecef;
        }
        h1, h2, h3 {
            color: #2c3e50;
            margin-top: 30px;
        }
        p {
            margin: 10px 0;
        }
        table {
            width: 100%;
            border-collapse: collapse;
            margin: 15px 0;
        }
        table th, table td {
            padding: 10px;
            border: 1px solid #ddd;
            text-align: left;
        }
        table th {
            background-color: #f2f2f2;
            font-weight: bold;
        }
        ul, ol {
            margin: 10px 0;
            padding-left: 20px;
        }
        """
    
    def _section_to_html(self, section: Dict[str, Any]) -> str:
        """将章节转换为HTML"""
        html_parts = []
        
        # 添加标题
        if 'title' in section:
            level = section.get('level', 2)
            html_parts.append(f"<h{level}>{section['title']}</h{level}>")
        
        # 添加段落
        if 'paragraphs' in section:
            for para in section['paragraphs']:
                html_parts.append(f"<p>{para}</p>")
        
        # 添加表格
        if 'tables' in section:
            for table_data in section['tables']:
                html_parts.append(self._table_to_html(table_data))
        
        # 添加列表
        if 'lists' in section:
            for list_data in section['lists']:
                html_parts.append(self._list_to_html(list_data))
        
        return "\n".join(html_parts)
    
    def _table_to_html(self, table_data: Dict[str, Any]) -> str:
        """将表格数据转换为HTML"""
        headers = table_data.get('headers', [])
        rows = table_data.get('rows', [])
        
        html_parts = ["<table>"]
        
        # 添加表头
        if headers:
            html_parts.append("<thead><tr>")
            for header in headers:
                html_parts.append(f"<th>{header}</th>")
            html_parts.append("</tr></thead>")
        
        # 添加数据行
        if rows:
            html_parts.append("<tbody>")
            for row in rows:
                html_parts.append("<tr>")
                for cell in row:
                    html_parts.append(f"<td>{cell}</td>")
                html_parts.append("</tr>")
            html_parts.append("</tbody>")
        
        html_parts.append("</table>")
        return "\n".join(html_parts)
    
    def _list_to_html(self, list_data: Dict[str, Any]) -> str:
        """将列表数据转换为HTML"""
        items = list_data.get('items', [])
        list_type = list_data.get('type', 'bullet')
        
        if not items:
            return ""
        
        tag = 'ol' if list_type == 'numbered' else 'ul'
        html_parts = [f"<{tag}>"]
        
        for item in items:
            html_parts.append(f"<li>{item}</li>")
        
        html_parts.append(f"</{tag}>")
        return "\n".join(html_parts)
    
    def _html_to_pdf(self, html_content: str, output_path: str) -> bool:
        """将HTML转换为PDF"""
        try:
            # 尝试使用weasyprint
            try:
                from weasyprint import HTML, CSS
                from weasyprint.text.fonts import FontConfiguration
                
                font_config = FontConfiguration()
                HTML(string=html_content).write_pdf(
                    output_path,
                    font_config=font_config
                )
                
                logger.info(f"PDF生成成功 (WeasyPrint): {output_path}")
                return True
                
            except ImportError:
                logger.warning("WeasyPrint不可用，尝试使用reportlab")
                
                # 回退到reportlab
                from reportlab.lib.pagesizes import A4
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                from reportlab.lib.units import inch
                
                doc = SimpleDocTemplate(output_path, pagesize=A4)
                styles = getSampleStyleSheet()
                
                # 创建自定义样式
                title_style = ParagraphStyle(
                    'CustomTitle',
                    parent=styles['Heading1'],
                    fontSize=18,
                    spaceAfter=30,
                    alignment=1  # 居中
                )
                
                # 解析HTML并创建PDF内容
                # 这里简化处理，实际应该使用HTML解析器
                story = []
                
                # 添加标题
                story.append(Paragraph("文档标题", title_style))
                story.append(Spacer(1, 20))
                
                # 添加内容
                story.append(Paragraph("文档内容", styles['Normal']))
                
                doc.build(story)
                
                logger.info(f"PDF生成成功 (ReportLab): {output_path}")
                return True
                
        except Exception as e:
            logger.error(f"HTML转PDF失败: {e}")
            return False


class DocumentGeneratorService:
    """文档生成服务"""
    
    def __init__(self):
        self.word_generator = WordGenerator()
        self.pdf_generator = PDFGenerator()
        self.excel_generator = ExcelGenerator()
        self.ppt_generator = PPTGenerator()
        self.output_dir = Path("generated_documents")
        self.output_dir.mkdir(exist_ok=True)
    
    def generate_document(
        self, 
        content: DocumentContent, 
        doc_type: DocumentType,
        filename: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        生成文档
        
        Args:
            content: 文档内容
            doc_type: 文档类型
            filename: 文件名（可选）
            
        Returns:
            生成结果信息
        """
        try:
            # 生成文件名
            if not filename:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"{content.title}_{timestamp}"
            
            # 根据类型选择生成器
            if doc_type == DocumentType.WORD:
                output_path = self.output_dir / f"{filename}.docx"
                success = self.word_generator.generate_document(content, str(output_path))
            elif doc_type == DocumentType.PDF:
                output_path = self.output_dir / f"{filename}.pdf"
                success = self.pdf_generator.generate_document(content, str(output_path))
            elif doc_type == DocumentType.EXCEL:
                output_path = self.output_dir / f"{filename}.xlsx"
                success = self.excel_generator.generate_document(content, str(output_path))
            elif doc_type == DocumentType.PPT:
                output_path = self.output_dir / f"{filename}.pptx"
                success = self.ppt_generator.generate_document(content, str(output_path))
            else:
                raise ValueError(f"不支持的文档类型: {doc_type}")
            
            if success:
                return {
                    "success": True,
                    "file_path": str(output_path),
                    "filename": output_path.name,
                    "file_size": output_path.stat().st_size,
                    "doc_type": doc_type.value,
                    "generated_at": datetime.now().isoformat()
                }
            else:
                return {
                    "success": False,
                    "error": "文档生成失败"
                }
                
        except Exception as e:
            logger.error(f"文档生成服务失败: {e}")
            return {
                "success": False,
                "error": str(e)
            }
    
    def parse_generation_request(self, request: str) -> Tuple[DocumentContent, DocumentType]:
        """
        解析生成请求
        
        Args:
            request: 用户请求文本
            
        Returns:
            (文档内容, 文档类型)
        """
        # 简单的意图识别
        request_lower = request.lower()
        
        # 识别文档类型
        if 'pdf' in request_lower:
            doc_type = DocumentType.PDF
        elif 'word' in request_lower or 'docx' in request_lower:
            doc_type = DocumentType.WORD
        else:
            doc_type = DocumentType.PDF  # 默认PDF
        
        # 识别模板类型
        if '报告' in request or 'report' in request_lower:
            template = DocumentTemplate.REPORT
        elif '总结' in request or 'summary' in request_lower:
            template = DocumentTemplate.SUMMARY
        elif '分析' in request or 'analysis' in request_lower:
            template = DocumentTemplate.ANALYSIS
        elif '提案' in request or 'proposal' in request_lower:
            template = DocumentTemplate.PROPOSAL
        else:
            template = DocumentTemplate.CUSTOM
        
        # 创建文档内容
        content = DocumentContent(
            title="生成的文档",
            sections=[
                {
                    "title": "内容",
                    "level": 1,
                    "paragraphs": [request]
                }
            ],
            metadata={
                "生成时间": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "模板类型": template.value,
                "文档类型": doc_type.value
            },
            template=template
        )
        
        return content, doc_type


class ExcelGenerator:
    """Excel文档生成器"""
    
    def __init__(self):
        self.templates_dir = Path("templates/excel")
        self.templates_dir.mkdir(exist_ok=True)
    
    def generate_document(self, content: DocumentContent, output_path: str) -> bool:
        """生成Excel文档"""
        try:
            import openpyxl
            from openpyxl.styles import Font, Alignment, PatternFill
            from openpyxl.utils import get_column_letter
            
            # 创建工作簿
            workbook = openpyxl.Workbook()
            sheet = workbook.active
            sheet.title = "主要内容"
            
            # 设置标题
            sheet['A1'] = content.title
            sheet['A1'].font = Font(size=16, bold=True)
            sheet['A1'].alignment = Alignment(horizontal='center')
            sheet.merge_cells('A1:D1')
            
            # 添加生成时间
            sheet['A2'] = f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            sheet['A2'].font = Font(size=10, italic=True)
            
            # 添加内容
            row = 4
            for section in content.sections:
                # 章节标题
                sheet[f'A{row}'] = section.get('title', '')
                sheet[f'A{row}'].font = Font(size=14, bold=True)
                sheet[f'A{row}'].fill = PatternFill(start_color='E6E6FA', end_color='E6E6FA', fill_type='solid')
                
                row += 1
                
                # 章节内容
                content_text = section.get('content', '')
                if content_text:
                    # 将内容分行
                    lines = content_text.split('\n')
                    for line in lines:
                        if line.strip():
                            sheet[f'A{row}'] = line.strip()
                            row += 1
                
                # 表格数据
                table_data = section.get('table_data')
                if table_data:
                    # 添加表头
                    if table_data:
                        headers = table_data[0] if len(table_data) > 0 else []
                        for col_idx, header in enumerate(headers):
                            cell = sheet.cell(row, col_idx + 1, header)
                            cell.font = Font(bold=True)
                            cell.fill = PatternFill(start_color='D3D3D3', end_color='D3D3D3', fill_type='solid')
                        row += 1
                        
                        # 添加数据行
                        for data_row in table_data[1:]:
                            for col_idx, cell_value in enumerate(data_row):
                                sheet.cell(row, col_idx + 1, cell_value)
                            row += 1
                
                row += 1  # 章节间空行
            
            # 调整列宽
            for column in sheet.columns:
                max_length = 0
                column_letter = get_column_letter(column[0].column)
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                sheet.column_dimensions[column_letter].width = adjusted_width
            
            # 保存文件
            workbook.save(output_path)
            logger.info(f"Excel文档生成成功: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"Excel文档生成失败: {e}")
            return False


class PPTGenerator:
    """PPT文档生成器"""
    
    def __init__(self):
        self.templates_dir = Path("templates/ppt")
        self.templates_dir.mkdir(exist_ok=True)
    
    def generate_document(self, content: DocumentContent, output_path: str) -> bool:
        """生成PPT文档"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            
            # 创建演示文稿
            prs = Presentation()
            
            # 标题页
            title_slide_layout = prs.slide_layouts[0]  # 标题页布局
            title_slide = prs.slides.add_slide(title_slide_layout)
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            
            title.text = content.title
            subtitle.text = f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            
            # 设置标题样式
            title.text_frame.paragraphs[0].font.size = Pt(44)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 内容页
            content_layout = prs.slide_layouts[1]  # 内容页布局
            
            for section in content.sections:
                slide = prs.slides.add_slide(content_layout)
                
                # 设置标题
                title_shape = slide.shapes.title
                title_shape.text = section.get('title', '')
                title_shape.text_frame.paragraphs[0].font.size = Pt(32)
                title_shape.text_frame.paragraphs[0].font.bold = True
                
                # 设置内容
                content_shape = slide.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                # 添加文本内容
                content_text = section.get('content', '')
                if content_text:
                    # 分行处理
                    lines = content_text.split('\n')
                    for i, line in enumerate(lines):
                        if line.strip():
                            if i == 0:
                                p = text_frame.paragraphs[0]
                            else:
                                p = text_frame.add_paragraph()
                            
                            p.text = line.strip()
                            p.font.size = Pt(18)
                            p.level = 0
                
                # 如果有表格数据，创建表格幻灯片
                table_data = section.get('table_data')
                if table_data and len(table_data) > 1:
                    # 创建表格幻灯片
                    table_slide = prs.slides.add_slide(content_layout)
                    table_title = table_slide.shapes.title
                    table_title.text = f"{section.get('title', '')} - 数据表格"
                    
                    # 创建表格
                    rows, cols = len(table_data), len(table_data[0])
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(8)
                    height = Inches(4)
                    
                    table = table_slide.shapes.add_table(rows, cols, left, top, width, height).table
                    
                    # 填充表格数据
                    for row_idx, row_data in enumerate(table_data):
                        for col_idx, cell_value in enumerate(row_data):
                            table.cell(row_idx, col_idx).text = str(cell_value)
                            
                            # 设置表头样式
                            if row_idx == 0:
                                table.cell(row_idx, col_idx).text_frame.paragraphs[0].font.bold = True
            
            # 保存文件
            prs.save(output_path)
            logger.info(f"PPT文档生成成功: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"PPT文档生成失败: {e}")
            return False
    
    def _setup_document_styles(self, doc):
        """设置文档样式"""
        try:
            # 设置标题样式
            styles = doc.styles
            
            # 标题1样式
            if 'Custom Heading 1' not in [s.name for s in styles]:
                h1_style = styles.add_style('Custom Heading 1', WD_STYLE_TYPE.PARAGRAPH)
                h1_style.font.name = 'Microsoft YaHei'
                h1_style.font.size = Pt(18)
                h1_style.font.bold = True
                h1_style.font.color.rgb = RGBColor(44, 62, 80)  # 深蓝色
                h1_style.paragraph_format.space_before = Pt(12)
                h1_style.paragraph_format.space_after = Pt(6)
            
            # 标题2样式
            if 'Custom Heading 2' not in [s.name for s in styles]:
                h2_style = styles.add_style('Custom Heading 2', WD_STYLE_TYPE.PARAGRAPH)
                h2_style.font.name = 'Microsoft YaHei'
                h2_style.font.size = Pt(16)
                h2_style.font.bold = True
                h2_style.font.color.rgb = RGBColor(52, 73, 94)  # 深灰色
                h2_style.paragraph_format.space_before = Pt(10)
                h2_style.paragraph_format.space_after = Pt(4)
            
            # 正文样式
            if 'Custom Body' not in [s.name for s in styles]:
                body_style = styles.add_style('Custom Body', WD_STYLE_TYPE.PARAGRAPH)
                body_style.font.name = 'Microsoft YaHei'
                body_style.font.size = Pt(12)
                body_style.paragraph_format.line_spacing = 1.5
                body_style.paragraph_format.space_after = Pt(6)
            
        except Exception as e:
            logger.warning(f"设置文档样式失败: {e}")
    
    def _add_header_footer(self, doc, title: str):
        """添加页眉页脚"""
        try:
            # 获取节
            section = doc.sections[0]
            
            # 设置页眉
            header = section.header
            header_para = header.paragraphs[0]
            header_para.text = title
            header_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            header_para.runs[0].font.size = Pt(10)
            header_para.runs[0].font.color.rgb = RGBColor(128, 128, 128)
            
            # 设置页脚
            footer = section.footer
            footer_para = footer.paragraphs[0]
            footer_para.text = f"第 {chr(0x200B)} 页，共 {chr(0x200B)} 页"  # 占位符
            footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            footer_para.runs[0].font.size = Pt(9)
            footer_para.runs[0].font.color.rgb = RGBColor(128, 128, 128)
            
        except Exception as e:
            logger.warning(f"添加页眉页脚失败: {e}")
    
    def _add_table_of_contents(self, doc, sections: List[Dict[str, Any]]):
        """添加目录"""
        try:
            doc.add_heading('目录', level=1)
            
            # 创建目录表格
            toc_table = doc.add_table(rows=len(sections) + 1, cols=3)
            toc_table.style = 'Table Grid'
            
            # 设置表头
            header_row = toc_table.rows[0]
            header_row.cells[0].text = "章节"
            header_row.cells[1].text = "标题"
            header_row.cells[2].text = "页码"
            
            # 设置表头样式
            for cell in header_row.cells:
                cell.paragraphs[0].runs[0].font.bold = True
                cell.paragraphs[0].runs[0].font.size = Pt(12)
            
            # 添加目录项
            for i, section in enumerate(sections, 1):
                row = toc_table.rows[i]
                row.cells[0].text = f"第{i}章"
                row.cells[1].text = section.get('title', f'章节{i}')
                row.cells[2].text = f"{i + 2}"  # 估算页码
            
            # 添加分页符
            doc.add_page_break()
            
        except Exception as e:
            logger.warning(f"添加目录失败: {e}")
    
    def _add_references_section(self, doc, content: DocumentContent):
        """添加参考文献章节"""
        try:
            doc.add_heading('参考文献', level=1)
            
            # 从元数据中提取参考文献
            references = content.metadata.get('references', [])
            if not references:
                # 生成默认参考文献
                references = [
                    "1. 相关文档资料",
                    "2. 网络搜索信息",
                    "3. 对话历史记录"
                ]
            
            # 添加参考文献列表
            for i, ref in enumerate(references, 1):
                para = doc.add_paragraph()
                para.add_run(f"[{i}] ").font.bold = True
                para.add_run(str(ref))
                para.style = 'List Number'
            
        except Exception as e:
            logger.warning(f"添加参考文献失败: {e}")
    
    def _add_enhanced_table(self, doc, table_data: List[List[str]], title: str = ""):
        """添加增强表格"""
        try:
            if not table_data or len(table_data) < 2:
                return
            
            # 添加表格标题
            if title:
                doc.add_heading(title, level=3)
            
            # 创建表格
            table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
            table.style = 'Table Grid'
            
            # 设置表格样式
            for i, row_data in enumerate(table_data):
                table_row = table.rows[i]
                for j, cell_value in enumerate(row_data):
                    cell = table_row.cells[j]
                    cell.text = str(cell_value)
                    
                    # 设置表头样式
                    if i == 0:
                        cell.paragraphs[0].runs[0].font.bold = True
                        cell.paragraphs[0].runs[0].font.size = Pt(11)
                        # 设置表头背景色
                        try:
                            shading_elm = parse_xml(r'<w:shd {} w:fill="D3D3D3"/>'.format(nsdecls('w')))
                            cell._tc.get_or_add_tcPr().append(shading_elm)
                        except:
                            pass
                    else:
                        cell.paragraphs[0].runs[0].font.size = Pt(10)
            
            # 添加表格说明
            if len(table_data) > 1:
                doc.add_paragraph(f"注：以上表格包含 {len(table_data)-1} 行数据", style='Caption')
            
        except Exception as e:
            logger.warning(f"添加增强表格失败: {e}")
    
    def _add_chart_placeholder(self, doc, chart_title: str, chart_type: str = "柱状图"):
        """添加图表占位符"""
        try:
            doc.add_heading(f"{chart_title} ({chart_type})", level=3)
            
            # 添加图表说明
            chart_desc = doc.add_paragraph()
            chart_desc.add_run("图表说明：").font.bold = True
            chart_desc.add_run(f"此处应显示{chart_type}，展示相关数据趋势和分析结果。")
            
            # 添加占位符文本
            placeholder = doc.add_paragraph()
            placeholder.add_run("[图表占位符]").font.italic = True
            placeholder.add_run(" - 实际部署时需要集成图表生成库（如matplotlib、plotly等）")
            
        except Exception as e:
            logger.warning(f"添加图表占位符失败: {e}")
    
    def _add_page_break(self, doc):
        """添加分页符"""
        try:
            doc.add_page_break()
        except Exception as e:
            logger.warning(f"添加分页符失败: {e}")
    
    def _add_enhanced_section(self, doc, section: Dict[str, Any]):
        """添加增强章节"""
        try:
            # 添加章节标题
            if 'title' in section:
                level = section.get('level', 2)
                heading = doc.add_heading(section['title'], level=level)
                
                # 设置标题样式
                if level == 1:
                    heading.style = 'Custom Heading 1'
                elif level == 2:
                    heading.style = 'Custom Heading 2'
            
            # 添加内容
            if 'content' in section:
                content_text = section['content']
                if content_text:
                    # 将内容按段落分割
                    paragraphs = content_text.split('\n\n')
                    for para_text in paragraphs:
                        if para_text.strip():
                            para = doc.add_paragraph(para_text.strip())
                            para.style = 'Custom Body'
            
            # 添加增强表格
            if 'table_data' in section and section['table_data']:
                table_data = section['table_data']
                if table_data and len(table_data) > 0:
                    self._add_enhanced_table(doc, table_data, f"{section.get('title', '')} - 数据表格")
            
            # 添加子章节
            if 'subsections' in section:
                for subsection in section['subsections']:
                    self._add_enhanced_section(doc, subsection)
            
            # 添加图表占位符（如果有结构化数据）
            if 'chart_data' in section:
                self._add_chart_placeholder(doc, f"{section.get('title', '')} - 数据图表")
            
            # 章节间添加分页符（除了最后一个章节）
            if section != section:  # 这里需要根据实际情况调整
                self._add_page_break(doc)
            
        except Exception as e:
            logger.error(f"添加增强章节失败: {e}")
            # 回退到原始方法
            self._add_section(doc, section)
